import sys
from pptx import Presentation

def safe(text):
    return text.encode('ascii', 'replace').decode('ascii')

prs = Presentation('theme/reference.pptx')
lines = []
lines.append(f'Size: {prs.slide_width.inches:.4f} x {prs.slide_height.inches:.4f} inches')
lines.append(f'Slides: {len(prs.slides)}')
lines.append(f'Layouts: {len(prs.slide_layouts)}')
lines.append('')

for si, slide in enumerate(prs.slides):
    lines.append(f'=== SLIDE {si+1} | layout: "{safe(slide.slide_layout.name)}" ===')
    for shape in slide.shapes:
        l = round(shape.left / 914400, 4) if shape.left is not None else None
        t = round(shape.top / 914400, 4) if shape.top is not None else None
        w = round(shape.width / 914400, 4) if shape.width is not None else None
        h = round(shape.height / 914400, 4) if shape.height is not None else None
        is_ph = shape.is_placeholder
        ph_idx = shape.placeholder_format.idx if is_ph else None
        text_preview = ''
        if shape.has_text_frame and shape.text_frame.text:
            text_preview = safe(shape.text_frame.text[:80])
        lines.append(f'  ph={ph_idx} name="{safe(shape.name)}" l={l} t={t} w={w} h={h}')
        if text_preview:
            lines.append(f'    text: {text_preview}')
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs[:4]):
                for ri, run in enumerate(para.runs[:1]):
                    fs = run.font.size
                    bold = run.font.bold
                    color = None
                    try:
                        color = run.font.color.rgb
                    except Exception:
                        pass
                    lines.append(f'    para{pi} run{ri}: size={round(fs.pt,1) if fs else "inh"}pt bold={bold} color={color} text={safe(run.text[:40])}')
    lines.append('')

with open('reference_info.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(lines))

print('Done. Lines written:', len(lines))
