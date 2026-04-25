"""Diagnostic: inspect reference.pptx layouts and slides."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

ref = os.path.join('theme', 'reference.pptx')
prs = Presentation(ref)
W, H = prs.slide_width, prs.slide_height
print(f'Slide dimensions: {W} x {H} EMU ({W/914400:.1f}in x {H/914400:.1f}in)')
print()

print('=== SLIDE LAYOUTS ===')
for i, layout in enumerate(prs.slide_layouts):
    phs = [(p.placeholder_format.idx, p.name, int(p.placeholder_format.type)) for p in layout.placeholders]
    non_ph = [(s.name, str(s.shape_type), s.width, s.height) for s in layout.shapes if not s.is_placeholder]
    print(f'  Layout {i}: "{layout.name}"')
    for idx, name, ptype in phs:
        print(f'    PH idx={idx}  type={ptype}  name={name}')
    for name, stype, w, h in non_ph:
        area_pct = (w * h) / (W * H) * 100
        print(f'    SHAPE: {name}  type={stype}  area={area_pct:.1f}%')
    print()

print('=== REFERENCE SLIDES ===')
for i, slide in enumerate(prs.slides):
    phs = [(p.placeholder_format.idx, p.name, int(p.placeholder_format.type)) for p in slide.placeholders]
    shapes = []
    for s in slide.shapes:
        if s.is_placeholder:
            continue
        area_pct = ((s.width or 0) * (s.height or 0)) / (W * H) * 100
        shapes.append((s.name, str(s.shape_type), area_pct, s.left, s.top, s.width, s.height))
    title_text = ''
    for p in slide.placeholders:
        if p.placeholder_format.idx == 0:
            title_text = (p.text or '')[:60]
            break
    print(f'  Slide {i}: "{title_text}"')
    for idx, name, ptype in phs:
        print(f'    PH idx={idx}  type={ptype}  name={name}')
    for name, stype, area, l, t, w, h in shapes:
        print(f'    SHAPE: {name}  type={stype}  area={area:.1f}%  pos=({l},{t}) size=({w},{h})')
    print()
