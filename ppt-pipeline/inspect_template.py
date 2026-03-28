from pptx import Presentation

prs = Presentation('theme/chosen_template.pptx')
print(f'Size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches')
print(f'Template slides: {len(prs.slides)}')
print()
print('=== LAYOUTS ===')
for i, layout in enumerate(prs.slide_layouts):
    phs = [(ph.placeholder_format.idx, ph.name) for ph in layout.placeholders]
    print(f'[{i}] "{layout.name}"')
    for idx, name in phs:
        print(f'     ph{idx}: {name}')
print()
print('=== EXISTING SLIDES ===')
for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    phs = [(ph.placeholder_format.idx, ph.name, ph.text[:50] if ph.text else '') for ph in slide.placeholders]
    print(f'Slide {i+1} (layout: "{layout_name}")')
    for idx, name, text in phs:
        print(f'     ph{idx}: {name} => "{text}"')
