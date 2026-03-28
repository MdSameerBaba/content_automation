import os
import traceback
from flask import Flask, request, jsonify, send_from_directory, send_file
from pipeline.stage1_parser import parse_pdf
from pipeline.checkpoint import CheckpointManager
from pipeline.stage2_structurer import structure_slides
from pipeline.stage3_content import generate_content

app = Flask(__name__, static_folder='static')
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
checkpoint_mgr = CheckpointManager()


@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')


@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '' or not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'No selected PDF file'}), 400
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(filepath)
    result = parse_pdf(filepath)
    return jsonify(result)


@app.route('/pipeline/status/<filename>')
def pipeline_status(filename):
    stages = [
        'stage1_parsed', 'stage2_structured', 'stage3_content',
        'stage4_pptx', 'stage5_images', 'stage6_audio', 'stage7_video'
    ]
    status = {}
    for stage in stages:
        status[stage] = checkpoint_mgr.exists(stage, filename)
    return jsonify(status)


@app.route('/checkpoint/<stage>/<filename>')
def get_checkpoint(stage, filename):
    data = checkpoint_mgr.load(stage, filename)
    if data is None:
        return jsonify({'error': 'Checkpoint not found'}), 404
    return jsonify(data)


@app.route('/checkpoints/<path:filepath>')
def serve_checkpoint_file(filepath):
    """Serve files from checkpoints directory."""
    from flask import send_from_directory
    base_dir = checkpoint_mgr.base_dir
    return send_from_directory(base_dir, filepath)
    return jsonify(data)


# ── Stage 2: Structure ───────────────────────────────────────────────
@app.route('/pipeline/structure/<filename>', methods=['POST'])
def run_structure(filename):
    try:
        result = structure_slides(filename)
        return jsonify({'success': True, 'result': result})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


# ── Stage 3: Content / AI gap fill + speaker notes ──────────────────
@app.route('/pipeline/content/<filename>', methods=['POST'])
def run_content(filename):
    try:
        result = generate_content(filename)
        return jsonify({
            'success': True,
            'audit': result['audit'],
            'missing_slides': result['missing_slides_content'],
            'speaker_notes_count': len(result['speaker_notes']),
            'original_slides': result['original_slide_count'],
            'final_slides': result['final_slide_count']
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


# ── Stage 4: Build PPTX ─────────────────────────────────────────────
@app.route('/pipeline/build/<filename>', methods=['POST'])
def run_build(filename):
    try:
        from pipeline.stage4_builder import build_pptx
        result = build_pptx(filename)
        return jsonify({'success': True, 'result': result})
    except ImportError:
        return jsonify({'success': False, 'error': 'stage4_builder.py not created yet'}), 501
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


# ── Download Stage 4 PPTX ───────────────────────────────────────────
@app.route('/download-pptx/<filename>')
def download_pptx(filename):
    """Serve the AI-generated PPTX for user to download and edit."""
    pptx_file = os.path.join(checkpoint_mgr.base_dir, 'stage4_pptx', f'{filename}.pptx')
    if not os.path.exists(pptx_file):
        return jsonify({'error': 'PPTX not found. Run Stage 4 first.'}), 404
    return send_file(pptx_file, as_attachment=True,
                     download_name=f'{filename}_ai_generated.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


# ── Download Stage 7 Video ─────────────────────────────────────────
@app.route('/download-video/<filename>')
def download_video(filename):
    """Serve the final MP4 video for download."""
    video_file = os.path.join(checkpoint_mgr.base_dir, 'stage7_video', f'{filename}.mp4')
    if not os.path.exists(video_file):
        return jsonify({'error': 'Video not found. Run Stage 7 first.'}), 404
    return send_file(video_file, as_attachment=True,
                     download_name=f'{filename}_video.mp4',
                     mimetype='video/mp4')


# ── Upload Human-Revised PPTX ────────────────────────────────────────
@app.route('/upload-revised-pptx', methods=['POST'])
def upload_revised_pptx():
    """
    Accept a user-edited PPTX and save as the Stage 5 input.
    This overrides the AI-generated version for all downstream stages.
    Speaker notes in the uploaded PPTX become the narration script.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    filename = request.form.get('filename', '')
    if not filename:
        return jsonify({'error': 'filename parameter required'}), 400
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Only .pptx files accepted'}), 400

    # Save to stage5_input as the authoritative source for stages 5-7
    input_dir = os.path.join(checkpoint_mgr.base_dir, 'stage5_input')
    os.makedirs(input_dir, exist_ok=True)
    save_path = os.path.join(input_dir, f'{filename}.pptx')
    file.save(save_path)

    # Read back slide count and notes word count as confirmation
    try:
        from pptx import Presentation
        prs = Presentation(save_path)
        notes_words = 0
        for slide in prs.slides:
            try:
                notes_words += len(slide.notes_slide.notes_text_frame.text.split())
            except Exception:
                pass
        slide_count = len(prs.slides)
    except Exception as e:
        slide_count = -1
        notes_words = -1

    return jsonify({
        'success': True,
        'saved_path': save_path,
        'slides': slide_count,
        'narrator_words': notes_words,
        'message': f'Revised PPTX saved. {slide_count} slides, ~{notes_words} narrator words.'
    })


# ── Stage 5: Export to PNG images ───────────────────────────────────
@app.route('/pipeline/images/<filename>', methods=['POST'])
def run_images(filename):
    try:
        from pipeline.stage5_images import export_images
        result = export_images(filename)
        return jsonify({'success': True, 'result': result})
    except ImportError:
        return jsonify({'success': False, 'error': 'stage5_images.py not created yet'}), 501
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


# ── Voice Selection ─────────────────────────────────────────────────
@app.route('/voices', methods=['GET'])
def get_voices():
    try:
        from pipeline.stage6_audio import get_voices, KOKORO_VOICES
        return jsonify({'success': True, 'voices': KOKORO_VOICES})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/voices/preview', methods=['POST'])
def preview_voice():
    try:
        data = request.get_json()
        voice_id = data.get('voice_id')
        custom_text = data.get('text')
        
        from pipeline.stage6_audio import generate_preview
        result = generate_preview(voice_id, custom_text)
        
        if 'error' in result:
            return jsonify({'success': False, 'error': result['error']}), 400
        
        return jsonify({'success': True, 'result': result})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


# ── Stage 6: Generate audio + SRT ───────────────────────────────────
@app.route('/pipeline/audio/<filename>', methods=['POST'])
def run_audio(filename):
    try:
        from pipeline.stage6_audio import generate_audio
        data = request.get_json() or {}
        voice = data.get('voice')  # Get selected voice
        result = generate_audio(filename, voice=voice)
        return jsonify({'success': True, 'result': result})
    except ImportError:
        return jsonify({'success': False, 'error': 'stage6_audio.py not created yet'}), 501
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


# ── Stage 7: Assemble MP4 video ─────────────────────────────────────
@app.route('/pipeline/video/<filename>', methods=['POST'])
def run_video(filename):
    try:
        from pipeline.stage7_video import create_video
        result = create_video(filename)
        return jsonify({'success': True, 'result': result})
    except ImportError:
        return jsonify({'success': False, 'error': 'stage7_video.py not created yet'}), 501
    except Exception as e:
        return jsonify({'success': False, 'error': str(e), 'traceback': traceback.format_exc()}), 500


if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True)
