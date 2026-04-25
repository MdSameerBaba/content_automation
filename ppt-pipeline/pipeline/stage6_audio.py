"""Stage 6: Generate narration audio (MP3) from PPTX speaker notes.

This module supports:
- Multilingual narration via Sarvam
- Provider-aware routing and fallback chain
- English transcript artifacts (JSON + SRT)
- Protected-keyword retention for technical terms
"""

import base64
import json
import os
import re
import subprocess
import tempfile

import imageio_ffmpeg
import requests
from pptx import Presentation

from pipeline.checkpoint import CheckpointManager, is_cache_reuse_enabled
from pipeline.config import get_tts_config


KOKORO_VOICES = {
    'af_heart': {'name': 'Heart (Female)', 'gender': 'Female'},
    'af_sarah': {'name': 'Sarah (Female)', 'gender': 'Female'},
    'am_adam': {'name': 'Adam (Male)', 'gender': 'Male'},
    'am_michael': {'name': 'Michael (Male)', 'gender': 'Male'},
    'bf_emma': {'name': 'Emma (British Female)', 'gender': 'Female'},
    'bm_george': {'name': 'George (British Male)', 'gender': 'Male'},
    'af_nova': {'name': 'Nova (Female)', 'gender': 'Female'},
}

SARVAM_LANGUAGES = {
    'en-IN': 'English (India)',
    'hi-IN': 'Hindi',
    'bn-IN': 'Bengali',
    'ta-IN': 'Tamil',
    'te-IN': 'Telugu',
    'kn-IN': 'Kannada',
    'ml-IN': 'Malayalam',
    'mr-IN': 'Marathi',
    'gu-IN': 'Gujarati',
    'pa-IN': 'Punjabi',
    'od-IN': 'Odia',
}

SARVAM_SPEAKERS = [
    'shubh', 'aditya', 'ritu', 'priya', 'neha', 'rahul', 'pooja', 'rohan', 'simran', 'kavya',
    'amit', 'dev', 'ishita', 'shreya', 'ratan', 'varun', 'manan', 'sumit', 'roopa', 'kabir',
    'aayan', 'ashutosh', 'advait', 'anand', 'tanya', 'tarun', 'sunny', 'mani', 'gokul', 'vijay',
    'shruti', 'suhani', 'mohit', 'kavitha', 'rehan', 'soham', 'rupali',
]

TTS_CONFIG = get_tts_config()
DEFAULT_SILENCE_MS = TTS_CONFIG['silence_duration_ms']
PREVIEW_TEXT = TTS_CONFIG['preview_text']

SARVAM_API_URL = os.getenv('SARVAM_TTS_URL', 'https://api.sarvam.ai/text-to-speech').strip()
SARVAM_TRANSLATE_URL = os.getenv('SARVAM_TRANSLATE_URL', 'https://api.sarvam.ai/translate').strip()
SARVAM_TRANSLATE_MODEL = os.getenv('SARVAM_TRANSLATE_MODEL', 'sarvam-translate:v1').strip() or 'sarvam-translate:v1'
SARVAM_MODEL = os.getenv('SARVAM_TTS_MODEL', 'bulbul:v3').strip() or 'bulbul:v3'
SARVAM_DEFAULT_SPEAKER = os.getenv('SARVAM_DEFAULT_SPEAKER', 'shubh').strip().lower() or 'shubh'
SARVAM_DEFAULT_LANGUAGE = os.getenv('SARVAM_DEFAULT_LANGUAGE', 'en-IN').strip() or 'en-IN'
SARVAM_DEFAULT_PACE = float(os.getenv('SARVAM_TTS_PACE', '1.0') or 1.0)
SARVAM_DEFAULT_SAMPLE_RATE = int(os.getenv('SARVAM_TTS_SAMPLE_RATE', '24000') or 24000)
SARVAM_MAX_TEXT_CHARS = int(os.getenv('SARVAM_MAX_TEXT_CHARS', '2400') or 2400)

TRANSCRIPT_FORMAT_DEFAULT = os.getenv('STAGE6_TRANSCRIPT_FORMAT', 'both').strip().lower() or 'both'
TRANSCRIPT_LANGUAGE_MODE_DEFAULT = os.getenv('STAGE6_TRANSCRIPT_LANGUAGE_MODE', 'always_english').strip().lower() or 'always_english'
KEYWORD_POLICY_DEFAULT = os.getenv('TTS_KEYWORD_POLICY', 'keep_english').strip().lower() or 'keep_english'

_PROVIDER_DEFAULT_CHAIN = os.getenv('TTS_PROVIDER_ORDER', 'sarvam,kokoro,piper,gtts')
PROVIDER_ORDER = [p.strip().lower() for p in _PROVIDER_DEFAULT_CHAIN.split(',') if p.strip()]

_VALID_PROVIDERS = {'sarvam', 'kokoro', 'piper', 'gtts'}
PROVIDER_ORDER = [p for p in PROVIDER_ORDER if p in _VALID_PROVIDERS] or ['sarvam', 'kokoro', 'piper', 'gtts']

_DEFAULT_LANGUAGE_FROM_ENV = os.getenv('DEFAULT_TTS_LANGUAGE', SARVAM_DEFAULT_LANGUAGE).strip() or 'en-IN'
DEFAULT_LANGUAGE_CODE = _DEFAULT_LANGUAGE_FROM_ENV if _DEFAULT_LANGUAGE_FROM_ENV in SARVAM_LANGUAGES else 'en-IN'

DEFAULT_PROVIDER = os.getenv('DEFAULT_TTS_PROVIDER', 'auto').strip().lower() or 'auto'
if DEFAULT_PROVIDER not in _VALID_PROVIDERS and DEFAULT_PROVIDER != 'auto':
    DEFAULT_PROVIDER = 'auto'

ENGLISH_LANGUAGE_CODES = {'en-IN'}

# Map Sarvam language codes → gTTS lang codes (for fallback when Sarvam is unavailable)
_SARVAM_TO_GTTS_LANG = {
    'en-IN': 'en',
    'hi-IN': 'hi',
    'bn-IN': 'bn',
    'ta-IN': 'ta',
    'te-IN': 'te',
    'kn-IN': 'kn',
    'ml-IN': 'ml',
    'mr-IN': 'mr',
    'gu-IN': 'gu',
    'pa-IN': 'pa',
    'od-IN': 'or',  # Odia
}

_PROVIDER_LANGUAGE_SUPPORT = {
    'sarvam': set(SARVAM_LANGUAGES.keys()),
    'kokoro': set(ENGLISH_LANGUAGE_CODES),
    'piper': set(ENGLISH_LANGUAGE_CODES),
    # gTTS supports all Sarvam languages via _SARVAM_TO_GTTS_LANG mapping
    'gtts': set(_SARVAM_TO_GTTS_LANG.keys()),
}

checkpoint_mgr = CheckpointManager()
OUTPUT_DIR = os.path.join(checkpoint_mgr.base_dir, 'stage6_audio')


def _parse_keyword_list(raw):
    if raw is None:
        return []
    if isinstance(raw, list):
        items = raw
    else:
        items = str(raw).split(',')
    normalized = []
    seen = set()
    for item in items:
        kw = str(item).strip()
        if not kw:
            continue
        key = kw.lower()
        if key in seen:
            continue
        seen.add(key)
        normalized.append(kw)
    return normalized


def _default_protected_keywords():
    env_value = os.getenv('TTS_PROTECTED_KEYWORDS') or os.getenv('PROTECTED_KEYWORDS') or ''
    return _parse_keyword_list(env_value)


DEFAULT_PROTECTED_KEYWORDS = _default_protected_keywords()


def _resolve_transcript_formats(transcript_format):
    if transcript_format is None:
        raw_items = [TRANSCRIPT_FORMAT_DEFAULT]
    elif isinstance(transcript_format, list):
        raw_items = [str(x).strip().lower() for x in transcript_format]
    else:
        raw_items = [x.strip().lower() for x in str(transcript_format).split(',')]

    raw_items = [x for x in raw_items if x]
    if not raw_items:
        raw_items = ['both']

    if 'none' in raw_items and len(raw_items) == 1:
        return []

    formats = set()
    for item in raw_items:
        if item == 'both':
            formats.update({'json', 'srt'})
        elif item in {'json', 'srt'}:
            formats.add(item)

    if not formats and 'none' not in raw_items:
        formats.update({'json', 'srt'})

    ordered = []
    if 'json' in formats:
        ordered.append('json')
    if 'srt' in formats:
        ordered.append('srt')
    return ordered


def _normalize_transcript_language_mode(mode):
    candidate = (mode or TRANSCRIPT_LANGUAGE_MODE_DEFAULT).strip().lower()
    if candidate in {'always_english', 'english', 'en', 'en_only'}:
        return 'always_english'
    return 'always_english'


def _normalize_keyword_policy(policy):
    candidate = (policy or KEYWORD_POLICY_DEFAULT).strip().lower()
    if candidate in {'keep_english', 'translate_all'}:
        return candidate
    return 'keep_english'


def _build_voice_catalog():
    catalog = {}

    for voice_id, meta in KOKORO_VOICES.items():
        catalog[f'kokoro:{voice_id}'] = {
            'name': meta.get('name', voice_id),
            'gender': meta.get('gender', 'Unknown'),
            'provider': 'kokoro',
            'voice': voice_id,
            'languages': ['en-IN'],
        }

    for speaker in SARVAM_SPEAKERS:
        display = speaker.replace('_', ' ').title()
        catalog[f'sarvam:{speaker}'] = {
            'name': f'{display} (Sarvam)',
            'gender': 'Unknown',
            'provider': 'sarvam',
            'voice': speaker,
            'languages': list(SARVAM_LANGUAGES.keys()),
        }

    return catalog


VOICE_CATALOG = _build_voice_catalog()


def _normalize_language_code(language_code):
    code = (language_code or '').strip()
    if code in SARVAM_LANGUAGES:
        return code
    return DEFAULT_LANGUAGE_CODE


def _normalize_voice_id(voice_id):
    if not voice_id:
        return None
    raw = str(voice_id).strip()
    if raw in VOICE_CATALOG:
        return raw

    kokoro_key = f'kokoro:{raw}'
    if kokoro_key in VOICE_CATALOG:
        return kokoro_key

    sarvam_key = f'sarvam:{raw.lower()}'
    if sarvam_key in VOICE_CATALOG:
        return sarvam_key

    return None


def _default_voice_id():
    configured = _normalize_voice_id(os.getenv('DEFAULT_VOICE', TTS_CONFIG.get('default_voice', '')).strip())
    if configured:
        return configured
    return 'kokoro:af_heart'


DEFAULT_VOICE_ID = _default_voice_id()


def _split_voice_id(voice_id):
    normalized = _normalize_voice_id(voice_id) or DEFAULT_VOICE_ID
    provider, value = normalized.split(':', 1)
    return provider, value


def _provider_supports_language(provider, language_code):
    return language_code in _PROVIDER_LANGUAGE_SUPPORT.get(provider, set())


def _resolve_provider_chain(requested_provider=None, voice_provider=None, language_code='en-IN'):
    requested = (requested_provider or '').strip().lower()
    if requested not in _VALID_PROVIDERS:
        requested = ''

    route = {
        'requested_provider': requested or 'auto',
        'requested_language': language_code,
        'auto_routed': False,
        'reason': '',
        'effective_primary': '',
    }

    if language_code not in ENGLISH_LANGUAGE_CODES:
        # Non-English: Sarvam is primary; gTTS is a real language-aware fallback.
        chain = ['sarvam', 'gtts']

        if requested and requested not in ('sarvam', 'gtts'):
            route['auto_routed'] = True
            route['reason'] = f'provider_{requested}_does_not_support_{language_code}'
        elif not requested and voice_provider and voice_provider not in ('sarvam', 'gtts'):
            route['auto_routed'] = True
            route['reason'] = f'voice_provider_{voice_provider}_does_not_support_{language_code}'

        route['effective_primary'] = 'sarvam'
        return chain, route

    chain = []
    if requested:
        chain.append(requested)

    if not chain and voice_provider in _VALID_PROVIDERS:
        chain.append(voice_provider)

    for item in PROVIDER_ORDER:
        if item not in chain:
            chain.append(item)

    route['effective_primary'] = chain[0] if chain else 'sarvam'
    return chain, route


def _run_command(command, timeout):
    result = subprocess.run(command, capture_output=True, timeout=timeout, text=True)
    if result.returncode != 0:
        raise RuntimeError(f'Command failed ({result.returncode}): {" ".join(command)}\n{result.stderr[:300]}')
    return result


def _get_input_pptx(filename):
    revised = os.path.join(checkpoint_mgr.base_dir, 'stage5_input', f'{filename}.pptx')
    if os.path.exists(revised):
        return revised
    ai_gen = os.path.join(checkpoint_mgr.base_dir, 'stage4_pptx', f'{filename}.pptx')
    if os.path.exists(ai_gen):
        return ai_gen
    raise FileNotFoundError(f'No PPTX found for "{filename}".')


def _extract_notes_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    notes = {}
    for idx, slide in enumerate(prs.slides, 1):
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame if notes_slide is not None else None
            text = (text_frame.text or '').strip() if text_frame is not None else ''
        except Exception:
            text = ''
        notes[idx] = text
    return notes


def _split_text_for_sarvam(text, max_chars=SARVAM_MAX_TEXT_CHARS):
    content = (text or '').strip()
    if not content:
        return []
    if len(content) <= max_chars:
        return [content]

    sentences = re.split(r'(?<=[.!?])\s+', content)
    chunks = []
    current = ''

    for sent in sentences:
        sent = sent.strip()
        if not sent:
            continue

        if len(sent) > max_chars:
            if current:
                chunks.append(current)
                current = ''
            for i in range(0, len(sent), max_chars):
                chunks.append(sent[i:i + max_chars])
            continue

        candidate = f'{current} {sent}'.strip() if current else sent
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = sent

    if current:
        chunks.append(current)

    return chunks


def _sarvam_request(text, speaker, language_code):
    api_key = (os.getenv('SARVAM_API_KEY') or '').strip()
    if not api_key:
        raise RuntimeError('SARVAM_API_KEY is not configured')

    payload = {
        'text': text,
        'target_language_code': language_code,
        'model': SARVAM_MODEL,
        'speaker': speaker,
        'output_audio_codec': 'mp3',
        'speech_sample_rate': SARVAM_DEFAULT_SAMPLE_RATE,
        'pace': SARVAM_DEFAULT_PACE,
    }
    headers = {
        'api-subscription-key': api_key,
        'Content-Type': 'application/json',
    }

    response = requests.post(SARVAM_API_URL, headers=headers, json=payload, timeout=60)
    if response.status_code >= 400:
        detail = ''
        try:
            detail = response.json().get('error', {}).get('message', '')
        except Exception:
            detail = response.text[:300]
        raise RuntimeError(f'Sarvam API error {response.status_code}: {detail}')

    data = response.json()
    audios = data.get('audios') or []
    if not audios:
        raise RuntimeError('Sarvam API returned no audio payload')

    encoded = audios[0]
    if not isinstance(encoded, str) or not encoded.strip():
        raise RuntimeError('Sarvam API returned invalid base64 audio')

    return base64.b64decode(encoded)


def _concat_mp3_chunks(chunk_paths, out_path):
    ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as concat_file:
        for path in chunk_paths:
            escaped = path.replace("'", "'\\''")
            concat_file.write(f"file '{escaped}'\n")
        concat_path = concat_file.name

    try:
        _run_command([
            ffmpeg_path,
            '-y',
            '-f', 'concat',
            '-safe', '0',
            '-i', concat_path,
            '-c', 'copy',
            out_path,
        ], timeout=40)
    finally:
        try:
            os.unlink(concat_path)
        except OSError:
            pass


def _sarvam_translate(text, target_language_code, protected_keywords=None):
    """
    Translate English text to target_language_code via Sarvam Translate API.
    Protected keywords are temporarily replaced with placeholders before
    translation and restored afterwards so they remain in English script.

    Returns (translated_text, was_translation_attempted).
    Falls back to original text if translation fails.
    """
    api_key = (os.getenv('SARVAM_API_KEY') or '').strip()
    if not api_key:
        return text, False

    if target_language_code in ENGLISH_LANGUAGE_CODES:
        # No translation needed for English.
        return text, False

    clean = (text or '').strip()
    if not clean:
        return clean, False

    # ── Step 1: mask protected keywords with unique placeholders ──────────
    masked = clean
    placeholder_map = {}  # placeholder → original keyword
    keywords_sorted = sorted(protected_keywords or [], key=len, reverse=True)
    for idx, kw in enumerate(keywords_sorted):
        placeholder = f'__KW{idx}__'
        pattern = re.compile(rf'(?<!\w){re.escape(kw)}(?!\w)', flags=re.IGNORECASE)
        if pattern.search(masked):
            masked = pattern.sub(placeholder, masked)
            placeholder_map[placeholder] = kw

    # ── Step 2: call Sarvam Translate ─────────────────────────────────────
    payload = {
        'input': masked,
        'source_language_code': 'en-IN',
        'target_language_code': target_language_code,
        'model': SARVAM_TRANSLATE_MODEL,
    }
    headers = {
        'api-subscription-key': api_key,
        'Content-Type': 'application/json',
    }
    try:
        response = requests.post(SARVAM_TRANSLATE_URL, headers=headers, json=payload, timeout=30)
        response.raise_for_status()
        data = response.json()
        translated = (data.get('translated_text') or '').strip()
        if not translated:
            print(f'    Sarvam Translate returned empty for lang {target_language_code}, using original')
            return clean, True
    except Exception as ex:
        print(f'    Sarvam Translate failed ({ex}), falling back to English text for TTS')
        return clean, False

    # ── Step 3: restore English keywords in translated text ───────────────
    result = translated
    for placeholder, original_kw in placeholder_map.items():
        result = result.replace(placeholder, original_kw)

    print(f'    Translated {len(clean)} chars en-IN → {target_language_code} '
          f'({len(placeholder_map)} keywords kept in English)')
    return result, True


def _generate_sarvam(text, out_path, speaker=None, language_code=None, protected_keywords=None):
    speaker_id = (speaker or SARVAM_DEFAULT_SPEAKER).strip().lower()
    if speaker_id not in SARVAM_SPEAKERS:
        speaker_id = SARVAM_DEFAULT_SPEAKER

    lang = _normalize_language_code(language_code or SARVAM_DEFAULT_LANGUAGE)

    # ── Translate English speaker notes → target language ─────────────────
    # The speaker notes from Stage 3 are always in English. Sarvam TTS
    # expects input in the target language, so we translate first.
    narration_text = text
    if lang not in ENGLISH_LANGUAGE_CODES:
        narration_text, translated = _sarvam_translate(text, lang, protected_keywords=protected_keywords)
        if not translated:
            print(f'    Translation unavailable; sending English text directly to Sarvam TTS ({lang})')

    chunks = _split_text_for_sarvam(narration_text)
    if not chunks:
        return False

    try:
        if len(chunks) == 1:
            audio_bytes = _sarvam_request(chunks[0], speaker_id, lang)
            with open(out_path, 'wb') as f:
                f.write(audio_bytes)
            return os.path.exists(out_path) and os.path.getsize(out_path) > 100

        temp_files = []
        with tempfile.TemporaryDirectory() as temp_dir:
            for idx, chunk in enumerate(chunks, start=1):
                chunk_audio = _sarvam_request(chunk, speaker_id, lang)
                part_path = os.path.join(temp_dir, f'chunk_{idx:02d}.mp3')
                with open(part_path, 'wb') as f:
                    f.write(chunk_audio)
                temp_files.append(part_path)

            _concat_mp3_chunks(temp_files, out_path)
            return os.path.exists(out_path) and os.path.getsize(out_path) > 100
    except Exception as ex:
        print(f'    Sarvam failed: {ex}')
        return False


def _generate_gtts(text, out_path, language_code='en-IN'):
    try:
        from gtts import gTTS
        # Map Sarvam language code to gTTS language tag (e.g. hi-IN → hi)
        gtts_lang = _SARVAM_TO_GTTS_LANG.get(language_code, 'en')
        tts = gTTS(text=text, lang=gtts_lang, slow=False)
        tts.save(out_path)
        return os.path.exists(out_path) and os.path.getsize(out_path) > 100
    except Exception as ex:
        print(f'    gTTS failed: {ex}')
    return False


def _generate_kokoro(text, out_path, voice='af_heart'):
    wav_path = out_path.replace('.mp3', '.wav')
    try:
        from kokoro import KPipeline  # type: ignore
        import soundfile as sf  # type: ignore
        import numpy as np

        pipeline = KPipeline(lang_code='a', repo_id='hexgrad/Kokoro-82M')
        generator = pipeline(text, voice=voice, speed=1, split_pattern=r'\n+')

        audio_chunks = []
        for _, _, audio in generator:
            audio_chunks.append(audio)

        if audio_chunks:
            full_audio = np.concatenate(audio_chunks)
            sf.write(wav_path, full_audio, 24000)

            ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
            _run_command([
                ffmpeg_path, '-i', wav_path,
                '-acodec', 'libmp3lame', '-y', out_path,
            ], timeout=30)
            return os.path.exists(out_path)
    except Exception as ex:
        print(f'    Kokoro failed: {ex}')
    finally:
        if os.path.exists(wav_path):
            try:
                os.remove(wav_path)
            except OSError:
                pass
    return False


def _generate_piper(text, out_path):
    text_file = None
    wav_path = out_path.replace('.mp3', '.wav')
    try:
        text_file = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False)
        text_file.write(text)
        text_file.close()

        _run_command([
            'piper', '--text_file', text_file.name,
            '--output_file', wav_path,
        ], timeout=30)

        ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
        _run_command([
            ffmpeg_path, '-i', wav_path,
            '-acodec', 'libmp3lame', '-y', out_path,
        ], timeout=30)

        return True
    except Exception as ex:
        print(f'    Piper failed: {ex}')
    finally:
        if text_file is not None:
            try:
                os.unlink(text_file.name)
            except OSError:
                pass
        if os.path.exists(wav_path):
            try:
                os.remove(wav_path)
            except OSError:
                pass
    return False


def _generate_silence(out_path, duration_ms=DEFAULT_SILENCE_MS):
    temp_wav = out_path.replace('.mp3', '.wav')
    try:
        ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
        _run_command([
            ffmpeg_path,
            '-f', 'lavfi', '-i', 'anullsrc=r=44100:cl=mono',
            '-t', str(duration_ms / 1000), '-y', temp_wav,
        ], timeout=10)
        _run_command([
            ffmpeg_path, '-i', temp_wav, '-acodec', 'libmp3lame', '-y', out_path,
        ], timeout=10)
        if os.path.exists(temp_wav):
            os.remove(temp_wav)
    except Exception:
        with open(out_path, 'wb') as f:
            f.write(b'\xff\xfb\x90\x00' * 200)


def _synthesize_with_provider(provider, text, out_path, voice_id, language_code, protected_keywords=None):
    voice_provider, raw_voice = _split_voice_id(voice_id)

    if provider == 'sarvam':
        speaker = raw_voice if voice_provider == 'sarvam' else SARVAM_DEFAULT_SPEAKER
        return _generate_sarvam(text, out_path, speaker=speaker, language_code=language_code,
                                protected_keywords=protected_keywords)

    if provider == 'kokoro':
        kokoro_voice = raw_voice if voice_provider == 'kokoro' else 'af_heart'
        if kokoro_voice not in KOKORO_VOICES:
            kokoro_voice = 'af_heart'
        return _generate_kokoro(text, out_path, voice=kokoro_voice)

    if provider == 'piper':
        return _generate_piper(text, out_path)

    if provider == 'gtts':
        return _generate_gtts(text, out_path, language_code=language_code)

    return False


def _synthesize_with_chain(text, out_path, voice_id, language_code, requested_provider=None,
                           resolved_chain=None, routing_info=None, protected_keywords=None):
    if resolved_chain is None:
        voice_provider, _ = _split_voice_id(voice_id)
        chain, route = _resolve_provider_chain(
            requested_provider=requested_provider,
            voice_provider=voice_provider,
            language_code=language_code,
        )
    else:
        chain = list(resolved_chain)
        route = dict(routing_info or {})

    for provider in chain:
        success = _synthesize_with_provider(
            provider, text, out_path, voice_id, language_code,
            protected_keywords=protected_keywords,
        )
        if success:
            return provider, chain, route

    return 'silence', chain, route


def _apply_keyword_retention(text, protected_keywords, language_code, keyword_policy):
    clean = (text or '').strip()
    if not clean:
        return '', []

    if keyword_policy != 'keep_english':
        return clean, []

    if not protected_keywords:
        return clean, []

    processed = clean
    matched = []

    for keyword in sorted(protected_keywords, key=lambda x: len(x), reverse=True):
        pattern = re.compile(rf'(?<!\w){re.escape(keyword)}(?!\w)', flags=re.IGNORECASE)
        if pattern.search(processed):
            processed = pattern.sub(keyword, processed)
            matched.append(keyword)

    # For English mode this acts as case-normalization of configured technical terms.
    # For non-English mode this preserves selected terms in English script.
    return processed, sorted(set(matched))


def _get_audio_duration(audio_path, default_seconds=3.0):
    try:
        try:
            # MoviePy v2.x
            from moviepy import AudioFileClip
        except ImportError:
            # MoviePy v1.x fallback
            from moviepy.audio.io.AudioFileClip import AudioFileClip
        clip = AudioFileClip(audio_path)
        duration = float(clip.duration or default_seconds)
        clip.close()
        return max(0.1, duration)
    except Exception:
        return default_seconds


def _format_srt_time(seconds):
    total_ms = int(round(max(0.0, float(seconds)) * 1000))
    hrs = total_ms // 3600000
    mins = (total_ms % 3600000) // 60000
    secs = (total_ms % 60000) // 1000
    ms = total_ms % 1000
    return f'{hrs:02d}:{mins:02d}:{secs:02d},{ms:03d}'


def _write_transcript_json(path, transcript_payload):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(transcript_payload, f, ensure_ascii=False, indent=2)


def _write_transcript_srt(path, segments):
    lines = []
    for idx, seg in enumerate(segments, start=1):
        lines.append(str(idx))
        lines.append(f"{_format_srt_time(seg['start'])} --> {_format_srt_time(seg['end'])}")
        lines.append(seg['text'] if seg['text'] else '[No narration]')
        lines.append('')

    with open(path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines).strip() + '\n')


def _generate_transcripts(
    transcript_by_slide,
    audio_files,
    out_dir,
    transcript_formats,
    transcript_language_mode,
    narration_language_code,
    protected_keywords,
    matched_keywords,
):
    result = {
        'language_mode': transcript_language_mode,
        'transcript_language_code': 'en-IN',
        'narration_language_code': narration_language_code,
        'formats': list(transcript_formats),
        'files': {},
        'segments': 0,
        'protected_keywords': protected_keywords,
        'matched_keywords': matched_keywords,
    }

    if not transcript_formats:
        return result

    os.makedirs(out_dir, exist_ok=True)

    slide_items = []
    for slide, text in sorted(transcript_by_slide.items(), key=lambda kv: int(kv[0])):
        slide_items.append({'slide': int(slide), 'text': text})

    if 'json' in transcript_formats:
        json_path = os.path.join(out_dir, 'transcript_en.json')
        payload = {
            'language_mode': transcript_language_mode,
            'transcript_language_code': 'en-IN',
            'narration_language_code': narration_language_code,
            'slides': slide_items,
            'protected_keywords': protected_keywords,
            'matched_keywords': matched_keywords,
        }
        _write_transcript_json(json_path, payload)
        result['files']['json'] = json_path

    if 'srt' in transcript_formats:
        segments = []
        cursor = 0.0
        by_slide = {int(item.get('slide', 0)): item for item in audio_files}

        for item in slide_items:
            slide_num = int(item['slide'])
            audio_item = by_slide.get(slide_num) or {}
            audio_path = audio_item.get('path')
            duration = _get_audio_duration(audio_path, default_seconds=DEFAULT_SILENCE_MS / 1000.0) if audio_path else (DEFAULT_SILENCE_MS / 1000.0)
            start = cursor
            end = cursor + duration
            cursor = end

            segments.append({
                'slide': slide_num,
                'start': start,
                'end': end,
                'text': item.get('text', ''),
            })

        srt_path = os.path.join(out_dir, 'transcript_en.srt')
        _write_transcript_srt(srt_path, segments)
        result['files']['srt'] = srt_path
        result['segments'] = len(segments)

    return result


def get_voices():
    return VOICE_CATALOG


def get_voice_catalog():
    return VOICE_CATALOG


def get_language_catalog():
    return SARVAM_LANGUAGES


def get_default_voice_id():
    return DEFAULT_VOICE_ID


def get_default_language_code():
    return DEFAULT_LANGUAGE_CODE


def get_default_provider():
    return DEFAULT_PROVIDER


def get_default_transcript_format():
    return TRANSCRIPT_FORMAT_DEFAULT


def get_default_transcript_language_mode():
    return TRANSCRIPT_LANGUAGE_MODE_DEFAULT


def get_default_keyword_policy():
    return KEYWORD_POLICY_DEFAULT


def get_default_protected_keywords():
    return DEFAULT_PROTECTED_KEYWORDS


def generate_preview(
    voice_id,
    preview_text=None,
    language_code=None,
    provider=None,
    keyword_policy=None,
    protected_keywords=None,
):
    selected_voice = _normalize_voice_id(voice_id) or DEFAULT_VOICE_ID
    selected_language = _normalize_language_code(language_code)
    selected_keyword_policy = _normalize_keyword_policy(keyword_policy)
    effective_keywords = _parse_keyword_list(protected_keywords) if protected_keywords is not None else list(DEFAULT_PROTECTED_KEYWORDS)

    preview_dir = os.path.join(OUTPUT_DIR, 'previews')
    os.makedirs(preview_dir, exist_ok=True)

    voice_provider, raw_voice = _split_voice_id(selected_voice)

    # Product policy: preview copy remains English regardless of selected narration language.
    text = (PREVIEW_TEXT or '').strip()
    if not text:
        return {'error': 'Preview text is empty'}

    retained_text, matched_keywords = _apply_keyword_retention(
        text,
        effective_keywords,
        selected_language,
        selected_keyword_policy,
    )

    requested_provider = (provider or '').strip().lower() or voice_provider
    if requested_provider not in _VALID_PROVIDERS:
        requested_provider = voice_provider

    provider_chain, provider_routing = _resolve_provider_chain(
        requested_provider=requested_provider,
        voice_provider=voice_provider,
        language_code=selected_language,
    )

    out_provider_for_name = provider_chain[0] if provider_chain else voice_provider
    out_path = os.path.join(preview_dir, f'preview_{out_provider_for_name}_{raw_voice}_{selected_language}.mp3')

    engine, chain, routing = _synthesize_with_chain(
        text=retained_text,
        out_path=out_path,
        voice_id=selected_voice,
        language_code=selected_language,
        requested_provider=requested_provider,
        resolved_chain=provider_chain,
        routing_info=provider_routing,
        # Preview always stays in English; no translation needed, but pass keywords for consistency.
        protected_keywords=effective_keywords,
    )

    if engine == 'silence' or not os.path.exists(out_path):
        return {'error': 'Failed to generate preview. Check Sarvam/Kokoro/Piper/gTTS dependencies and keys.'}

    warning = None
    if engine != chain[0]:
        warning = f'Primary provider "{chain[0]}" was unavailable; preview generated with "{engine}" fallback.'

    return {
        'voice_id': selected_voice,
        'voice_name': VOICE_CATALOG.get(selected_voice, {}).get('name', selected_voice),
        'provider_requested': requested_provider,
        'provider_used': engine,
        'provider_chain': chain,
        'provider_routing': routing,
        'language_code': selected_language,
        'engine': engine,
        'preview_text_policy': 'english_only_fixed',
        'keyword_policy': selected_keyword_policy,
        'protected_keywords': effective_keywords,
        'matched_keywords': matched_keywords,
        'preview_url': f'/checkpoints/stage6_audio/previews/preview_{out_provider_for_name}_{raw_voice}_{selected_language}.mp3',
        **({'warning': warning} if warning else {}),
    }


def generate_audio(
    filename,
    voice=None,
    language_code=None,
    provider=None,
    transcript_format=None,
    transcript_language_mode=None,
    keyword_policy=None,
    protected_keywords=None,
):
    selected_voice = _normalize_voice_id(voice) or DEFAULT_VOICE_ID
    selected_language = _normalize_language_code(language_code)

    requested_provider = (provider or '').strip().lower() if provider else ''
    if requested_provider not in _VALID_PROVIDERS:
        requested_provider = ''

    selected_keyword_policy = _normalize_keyword_policy(keyword_policy)
    effective_keywords = _parse_keyword_list(protected_keywords) if protected_keywords is not None else list(DEFAULT_PROTECTED_KEYWORDS)

    selected_transcript_formats = _resolve_transcript_formats(transcript_format)
    selected_transcript_language_mode = _normalize_transcript_language_mode(transcript_language_mode)

    voice_provider, _ = _split_voice_id(selected_voice)
    provider_chain, provider_routing = _resolve_provider_chain(
        requested_provider=requested_provider,
        voice_provider=voice_provider,
        language_code=selected_language,
    )

    if is_cache_reuse_enabled() and checkpoint_mgr.exists('stage6_audio', filename):
        cached = checkpoint_mgr.load('stage6_audio', filename)
        cache_ok = (
            cached
            and 'error' not in cached
            and cached.get('voice_id') == selected_voice
            and cached.get('language_code') == selected_language
            and (cached.get('provider_requested') or '') == requested_provider
            and (cached.get('keyword_policy') or KEYWORD_POLICY_DEFAULT) == selected_keyword_policy
            and (cached.get('protected_keywords') or []) == effective_keywords
            and sorted(cached.get('transcript_formats') or []) == sorted(selected_transcript_formats)
            and (cached.get('transcript_language_mode') or TRANSCRIPT_LANGUAGE_MODE_DEFAULT) == selected_transcript_language_mode
            and 'engine_summary' in cached
        )
        if cache_ok:
            print(f'Valid Stage 6 checkpoint for {filename} (voice: {selected_voice}, lang: {selected_language})')
            return cached

    pptx_path = _get_input_pptx(filename)
    notes = _extract_notes_from_pptx(pptx_path)
    out_dir = os.path.join(OUTPUT_DIR, filename)
    os.makedirs(out_dir, exist_ok=True)

    audio_files = []
    total_words = 0
    transcript_source_by_slide = {}
    matched_keywords_all = set()

    voice_meta = VOICE_CATALOG.get(selected_voice, {})
    requested_display = requested_provider or 'auto'
    print(
        f'  Voice: {voice_meta.get("name", selected_voice)} | '
        f'Language: {selected_language} | '
        f'Requested provider: {requested_display} | '
        f'Primary provider: {provider_chain[0] if provider_chain else "unknown"}'
    )

    for slide_idx, text in notes.items():
        out_path = os.path.join(out_dir, f'slide_{slide_idx:02d}.mp3')
        clean_text = (text or '').strip()
        engine_used = 'silence'
        narration_text = clean_text

        if clean_text:
            narration_text, matched_keywords = _apply_keyword_retention(
                clean_text,
                effective_keywords,
                selected_language,
                selected_keyword_policy,
            )
            matched_keywords_all.update(matched_keywords)
            transcript_source_by_slide[slide_idx] = narration_text

            print(f'  Slide {slide_idx}: generating audio ({len(narration_text.split())} words)...')
            engine_used, _, _ = _synthesize_with_chain(
                text=narration_text,
                out_path=out_path,
                voice_id=selected_voice,
                language_code=selected_language,
                requested_provider=requested_provider,
                resolved_chain=provider_chain,
                routing_info=provider_routing,
                protected_keywords=effective_keywords,
            )

            if engine_used == 'silence':
                print(f'  Slide {slide_idx}: all providers failed, generating silence...')
                _generate_silence(out_path)
            else:
                total_words += len(narration_text.split())
        else:
            transcript_source_by_slide[slide_idx] = ''
            print(f'  Slide {slide_idx}: no notes -> generating silence...')
            _generate_silence(out_path)

        audio_files.append({
            'slide': slide_idx,
            'path': out_path,
            'words': len(narration_text.split()) if narration_text else 0,
            'has_speech': engine_used != 'silence',
            'engine': engine_used,
            'language_code': selected_language,
        })

    transcripts_meta = _generate_transcripts(
        transcript_by_slide=transcript_source_by_slide,
        audio_files=audio_files,
        out_dir=out_dir,
        transcript_formats=selected_transcript_formats,
        transcript_language_mode=selected_transcript_language_mode,
        narration_language_code=selected_language,
        protected_keywords=effective_keywords,
        matched_keywords=sorted(matched_keywords_all),
    )

    engine_summary = {}
    for item in audio_files:
        engine = item.get('engine', 'unknown')
        engine_summary[engine] = engine_summary.get(engine, 0) + 1

    non_silence = {k: v for k, v in engine_summary.items() if k != 'silence'}
    primary_engine = max(non_silence, key=lambda k: non_silence[k]) if non_silence else 'silence'

    chain_primary = provider_chain[0] if provider_chain else ''
    fallback_used = (
        any(item.get('engine') == 'silence' and item.get('words', 0) > 0 for item in audio_files)
        or any(item.get('engine') not in {'silence', chain_primary} for item in audio_files)
        or (primary_engine not in {'silence', chain_primary})
    )

    result = {
        'filename': filename,
        'voice_id': selected_voice,
        'voice': voice_meta.get('name', selected_voice),
        'provider_requested': requested_provider,
        'provider_used': primary_engine,
        'provider_chain': provider_chain,
        'provider_routing': provider_routing,
        'language_code': selected_language,
        'slide_count': len(audio_files),
        'total_words': total_words,
        'output_dir': out_dir,
        'audio_files': audio_files,
        'engine_summary': engine_summary,
        'primary_engine': primary_engine,
        'fallback_used': fallback_used,
        'transcript_formats': selected_transcript_formats,
        'transcript_language_mode': selected_transcript_language_mode,
        'transcripts': transcripts_meta,
        'keyword_policy': selected_keyword_policy,
        'protected_keywords': effective_keywords,
        'matched_keywords': sorted(matched_keywords_all),
    }
    checkpoint_mgr.save('stage6_audio', filename, result)
    return result
