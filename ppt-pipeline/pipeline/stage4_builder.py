"""
Stage 4 (v4): Build themed PPTX using reference.pptx as base.

Fixes:
- XML-level paragraph replacement (no double-content from reference)
- clean_body_lines(): strips bullet chars, filters orphans, joins continuations
- trim_notes(): caps notes at 100 words for ≤10min video target
- Formatting preserved via pPr/rPr template extraction from first para
"""

import os
import re
import copy
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pipeline.checkpoint import CheckpointManager, is_cache_reuse_enabled

checkpoint_mgr = CheckpointManager()

_PIPELINE_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
REFERENCE_PATH = os.path.join(_PIPELINE_ROOT, 'theme', 'reference.pptx')
OUTPUT_DIR = os.path.join(checkpoint_mgr.base_dir, 'stage4_pptx')

BULLET_CHARS = '●•·○◦▪▸►▶–—*-'
BUILD_POLICY_VERSION = 15


# ── Text cleaning ─────────────────────────────────────────────────────

def clean_body_lines(lines):
    """
    Clean raw parsed lines for slide body content:
    1. Strip leading bullet symbols and whitespace
    2. Filter orphan lines (< 3 chars or pure numbers)
    3. Join continuation lines (lowercase-start short fragments) to previous
    Returns at most 8 lines.
    """
    cleaned = []
    for line in lines:
        # Strip leading bullet chars + spaces
        stripped = line.lstrip(BULLET_CHARS + ' ').strip()
        # Filter empty / too-short / pure-number orphans
        if len(stripped) < 3 or stripped.isdigit():
            continue
        cleaned.append(stripped)

    # Join continuation lines to the previous one
    if cleaned:
        joined = [cleaned[0]]
        for line in cleaned[1:]:
            prev = joined[-1]
            is_continuation = (
                line[0].islower()                            # starts lowercase
                or (len(line.split()) <= 3                  # very short fragment
                    and not prev.endswith(('.', '!', '?', ':')))
            )
            if is_continuation:
                joined[-1] = prev + ' ' + line
            else:
                joined.append(line)
        cleaned = joined

    return cleaned[:8]


def trim_notes(notes, max_words=100):
    """Cap speaker notes at max_words for ≤10-min video pacing."""
    if not notes:
        return notes
    words = notes.split()
    if len(words) <= max_words:
        return notes
    # Trim to last complete sentence within limit
    trimmed = ' '.join(words[:max_words])
    last_stop = max(trimmed.rfind('.'), trimmed.rfind('!'), trimmed.rfind('?'))
    if last_stop > len(trimmed) // 2:
        return trimmed[:last_stop + 1]
    return trimmed + '...'


def _clean_text(text):
    return re.sub(r'\s+', ' ', (text or '')).strip()


def _trim_title(title, max_words=14):
    title = _clean_text(title)
    if not title:
        return 'Untitled Slide'
    words = title.split()
    if len(words) <= max_words:
        return title
    return ' '.join(words[:max_words]).rstrip() + '...'


def _trim_bullets(bullets, max_items=7, max_words_per_line=18):
    cleaned = []
    for raw in bullets or []:
        line = _clean_text(raw)
        if not line:
            continue
        words = line.split()
        if len(words) > max_words_per_line:
            line = ' '.join(words[:max_words_per_line]).rstrip() + '...'
        cleaned.append(line)
        if len(cleaned) >= max_items:
            break
    return cleaned


# ── XML helpers ───────────────────────────────────────────────────────

def _get_format_templates(txBody):
    """Extract pPr and rPr from the first non-empty paragraph as formatting templates."""
    paras = txBody.findall(qn('a:p'))
    for first_p in paras:
        runs = first_p.findall(qn('a:r'))
        if not runs:
            continue
        pPr = first_p.find(qn('a:pPr'))
        rPr = runs[0].find(qn('a:rPr'))
        return (copy.deepcopy(pPr) if pPr is not None else None,
                copy.deepcopy(rPr) if rPr is not None else None)
    return None, None


def _xml_replace_text(txBody, text_lines):
    """
    Remove ALL existing <a:p> elements and replace with one per text line.
    Preserves run/paragraph formatting from the original first paragraph.
    """
    pPr_tmpl, rPr_tmpl = _get_format_templates(txBody)

    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    for line in (text_lines if text_lines else ['']):
        p_elem = etree.SubElement(txBody, qn('a:p'))
        if pPr_tmpl is not None:
            p_elem.insert(0, copy.deepcopy(pPr_tmpl))
        r_elem = etree.SubElement(p_elem, qn('a:r'))
        if rPr_tmpl is not None:
            r_elem.insert(0, copy.deepcopy(rPr_tmpl))
        t_elem = etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = line


def _set_ph_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            _xml_replace_text(ph.text_frame._txBody, [text])
            return True
    return False


def _set_ph_bullets(slide, idx, lines):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            _xml_replace_text(ph.text_frame._txBody, lines if lines else [''])
            return True
    return False


def _resolve_body_placeholder_idx(slide, slide_type):
    """
    Resolve the body placeholder index based on slide type.
    Returns the placeholder index to use for body content.
    """
    body_indices = [1, 2]
    for idx in body_indices:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return idx
    return None


def _set_notes(slide, text):
    if not text:
        return
    try:
        tf = slide.notes_slide.notes_text_frame
        _xml_replace_text(tf._txBody, [text])
    except Exception as e:
        print(f'    warn: notes failed: {e}')


def _find_layout(prs, preferred_names):
    for name in preferred_names:
        for layout in prs.slide_layouts:
            if name.lower() in layout.name.lower():
                return layout
    return prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]


def _remove_content_area_pictures(slide, slide_type):
    """
    Remove ALL non-placeholder, non-corner-logo picture shapes from a reference slide.

    Called for every slide type before we write new content or insert our own images.
    This ensures reference template images never bleed through under new content.
    Brand logos (small, top-right corner) are preserved.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        try:
            prs = slide.part.package.presentation
            slide_width = prs.slide_width
            slide_height = prs.slide_height
        except Exception:
            slide_width, slide_height = 9144000, 6858000  # 10" x 7.56" standard

        logo_area_limit = slide_width * slide_height * 0.06  # < 6% = small corner logo

        to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            shape_area = (shape.width or 0) * (shape.height or 0)
            # Keep corner logos: small AND in the right 25% AND top 25%
            is_corner_logo = (
                shape_area < logo_area_limit
                and (shape.left or 0) > slide_width * 0.75
                and (shape.top or 0) < slide_height * 0.25
            )
            if not is_corner_logo:
                to_remove.append(shape)

        for shape in to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        if to_remove:
            print(f'    Cleaned {len(to_remove)} stale reference picture(s) ({slide_type})')
        return len(to_remove)
    except Exception as ex:
        print(f'    warn: _remove_content_area_pictures failed: {ex}')
        return 0


def _insert_source_image(slide, image_path, slide_type):
    """
    Insert the Stage 1 source PNG into the slide at the correct position.

    - diagram: image fills the full content area (below title, full width)
    - image_content: image fills the right half of the content area

    Silently skips if image_path is missing or the file does not exist.
    """
    if not image_path or not os.path.exists(image_path):
        return False
    if slide_type not in ('diagram', 'image_content'):
        return False

    try:
        from pptx.util import Emu

        try:
            prs = slide.part.package.presentation
            W = int(prs.slide_width)
            H = int(prs.slide_height)
        except Exception:
            W, H = 9144000, 6858000

        MARGIN = int(W * 0.03)          # 3% margin
        TITLE_H = int(H * 0.22)         # title occupies ~22% of slide height
        CONTENT_TOP = TITLE_H + MARGIN
        CONTENT_H = H - CONTENT_TOP - MARGIN

        if slide_type == 'diagram':
            # Full-width content area
            left = MARGIN
            top = CONTENT_TOP
            width = W - 2 * MARGIN
            height = CONTENT_H
        else:  # image_content: right half
            left = W // 2 + MARGIN
            top = CONTENT_TOP
            width = W // 2 - 2 * MARGIN
            height = CONTENT_H

        slide.shapes.add_picture(image_path, Emu(left), Emu(top), Emu(width), Emu(height))
        print(f'    Inserted source image ({slide_type}): {os.path.basename(image_path)}')
        return True
    except Exception as ex:
        print(f'    warn: _insert_source_image failed: {ex}')
        return False


# ── Layout mapping ────────────────────────────────────────────────────
# Based on diagnostic of reference.pptx:
#
#   Layout 0  SECTION_HEADER              title(0), footer(12)  — decorative group
#   Layout 1  TITLE                       title(0), subtitle(1), footer(12)
#   Layout 2  TITLE_AND_BODY              title(0), body(1), footer(12) — decorative group 24%
#   Layout 3  SECTION_TITLE_AND_DESCRIPTION  title(0), sub(1), body(2), footer(12) — 50% blue rect!
#   Layout 4  TITLE_AND_TWO_COLUMNS       title(0), body(1), body(2), footer(12) — CLEAN
#   Layout 5  TITLE_ONLY                  title(0), footer(12) — CLEAN
#   Layout 6  ONE_COLUMN_TEXT             title(0), body(1), footer(12) — CLEAN
#   Layout 7  MAIN_POINT                  title(0), footer(12) — decorative group
#   Layout 8  CAPTION_ONLY               body(1), footer(12)
#   Layout 9  BIG_NUMBER                  title(0), body(1), footer(12) — decorative group
#   Layout 10 BLANK                       footer(12)
#
# Strategy: only use CLEAN layouts (no decorative shapes that bleed through).
#   title         → Layout 1  TITLE
#   content/agenda/conclusion → Layout 6  ONE_COLUMN_TEXT (clean title + full-width body)
#   image_content → Layout 4  TITLE_AND_TWO_COLUMNS (left body=text, right body=image)
#   diagram       → Layout 5  TITLE_ONLY (title + we insert image ourselves)


def _pick_layout(prs, slide_type):
    """Select the correct clean layout for each slide type."""
    layout_map = {
        'title':         ['TITLE'],
        'content':       ['ONE_COLUMN_TEXT', 'TITLE_AND_BODY'],
        'agenda':        ['ONE_COLUMN_TEXT', 'TITLE_AND_BODY'],
        'conclusion':    ['ONE_COLUMN_TEXT', 'TITLE_AND_BODY'],
        'image_content': ['TITLE_AND_TWO_COLUMNS'],
        'diagram':       ['TITLE_ONLY'],
    }
    preferred = layout_map.get(slide_type, ['ONE_COLUMN_TEXT', 'TITLE_AND_BODY'])
    return _find_layout(prs, preferred)


def _insert_source_image(slide, image_path, slide_type):
    """
    Insert the Stage 1 source PNG into the slide at the correct position.

    - diagram: image fills the full content area below the title
    - image_content: image fills the right body column area

    Silently skips if the file does not exist.
    """
    if not image_path or not os.path.exists(image_path):
        return False
    if slide_type not in ('diagram', 'image_content'):
        return False

    try:
        from pptx.util import Emu

        try:
            prs_obj = slide.part.package.presentation
            W = int(prs_obj.slide_width)
            H = int(prs_obj.slide_height)
        except Exception:
            W, H = 9144000, 5143500

        if slide_type == 'diagram':
            # Full content area below the title row.
            # Title occupies roughly top 18%, leave a small gap.
            left = int(W * 0.05)
            top = int(H * 0.20)
            width = int(W * 0.90)
            height = int(H * 0.75)
        else:
            # image_content: right half of the content area.
            # Matches TITLE_AND_TWO_COLUMNS body(2) position.
            left = int(W * 0.50) + int(W * 0.02)
            top = int(H * 0.22)
            width = int(W * 0.45)
            height = int(H * 0.70)

        slide.shapes.add_picture(image_path, Emu(left), Emu(top), Emu(width), Emu(height))
        print(f'    Inserted source image ({slide_type}): {os.path.basename(image_path)}')
        return True
    except Exception as ex:
        print(f'    warn: _insert_source_image failed: {ex}')
        return False


# ── Main builder ──────────────────────────────────────────────────────

def build_pptx(filename):
    """Stage 4: build themed PPTX from stage3 manifest using reference.pptx."""

    parsed = checkpoint_mgr.load('stage1_parsed', filename)
    stage3 = checkpoint_mgr.load('stage3_content', filename)

    if not parsed:
        raise Exception('Stage 1 checkpoint missing.')
    if not stage3:
        raise Exception('Stage 3 checkpoint missing.')

    source_blueprint_version = int(stage3.get('typed_blueprint_version', 0) or 0)
    source_notes_policy_version = int(stage3.get('notes_policy_version', 0) or 0)

    if is_cache_reuse_enabled() and checkpoint_mgr.exists('stage4_pptx', filename):
        cached = checkpoint_mgr.load('stage4_pptx', filename)
        if (
            cached
            and 'error' not in cached
            and cached.get('build_policy_version') == BUILD_POLICY_VERSION
            and int(cached.get('source_stage3_blueprint_version', 0) or 0) == source_blueprint_version
            and int(cached.get('source_stage3_notes_policy_version', 0) or 0) == source_notes_policy_version
        ):
            print(f'Valid stage 4 checkpoint for {filename}')
            return cached

    ordered = _build_ordered_manifest(parsed, stage3)

    print(f'Building {len(ordered)} slides...')

    # ── Open reference — use its LAYOUTS only ─────────────────────────
    if not os.path.exists(REFERENCE_PATH):
        raise FileNotFoundError(
            f'Reference template not found: {REFERENCE_PATH}. '
            'Expected file at ppt-pipeline/theme/reference.pptx'
        )
    prs = Presentation(REFERENCE_PATH)

    # Remove ALL reference slides. We only want the slide layouts (theme,
    # fonts, colors). Individual reference slides contain position-specific
    # images, decorative shapes, and content that bleeds through if reused.
    ref_count = len(prs.slides)
    for idx in range(ref_count - 1, -1, -1):
        _remove_slide(prs, idx)
    print(f'Cleared {ref_count} reference slides; building {len(ordered)} from clean layouts')

    slide_manifest = []

    for i, item in enumerate(ordered):
        title = item['title']
        bullets = item.get('bullets', [])
        notes = item.get('notes', '')
        slide_type = item.get('type', 'content')
        embed_image = bool(item.get('embed_source_image', False))
        image_path = item.get('image_path') or ''
        has_source_image = embed_image and image_path and os.path.exists(image_path)

        # Create a fresh slide from the correct clean layout
        layout = _pick_layout(prs, slide_type)
        slide = prs.slides.add_slide(layout)

        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        print(f'  [{i+1}] {slide_type} "{title[:50]}" | layout={layout.name} | phs={ph_indices} | {len(bullets)} bullets')

        effective_bullets = _trim_bullets(bullets)

        # ── Set title (placeholder idx 0) ─────────────────────────────
        _set_ph_text(slide, 0, title)

        # ── Type-specific content placement ───────────────────────────
        if slide_type == 'title':
            # TITLE layout has subtitle at idx 1
            _set_ph_text(slide, 1, '')

        elif slide_type == 'diagram':
            # TITLE_ONLY layout: no body placeholder, just insert image
            if has_source_image:
                _insert_source_image(slide, image_path, 'diagram')

        elif slide_type == 'image_content':
            # TITLE_AND_TWO_COLUMNS: body(1)=left text, body(2)=right text
            # Write bullets to left column (idx 1)
            _set_ph_bullets(slide, 1, effective_bullets if effective_bullets else [''])
            # Clear right column placeholder — the source image will cover it
            _set_ph_bullets(slide, 2, [''])
            if has_source_image:
                _insert_source_image(slide, image_path, 'image_content')

        else:
            # content / agenda / conclusion → ONE_COLUMN_TEXT: body at idx 1
            _set_ph_bullets(slide, 1, effective_bullets if effective_bullets else [''])

        _set_notes(slide, notes)

        slide_manifest.append({
            'slide_num': i + 1,
            'type': slide_type,
            'title': title,
            'bullets': effective_bullets,
            'notes': notes,
            'source_slide_num': item.get('source_slide_num'),
            'title_inferred': bool(item.get('title_inferred', False)),
            'embed_source_image': bool(item.get('embed_source_image', False)),
            'image_path': item.get('image_path') or '',
        })

    # Save the PPTX
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f'{filename}.pptx')
    prs.save(out_path)
    print(f'Saved PPTX to: {out_path}')

    # Calculate totals
    total_words = sum(len(item['notes'].split()) for item in ordered if item.get('notes'))
    est_mins = total_words / 130

    print(f'\nNarrator: ~{total_words} words -> ~{est_mins:.1f} min at 130 wpm')
    final_slide_count = len(prs.slides)
    print(f'Stage 4 complete: {final_slide_count} slides -> {out_path}')

    result = {
        'filename': filename,
        'output_path': out_path,
        'total_slides': final_slide_count,
        'narrator_words': total_words,
        'estimated_minutes': round(est_mins, 2),
        'build_policy_version': BUILD_POLICY_VERSION,
        'source_stage3_blueprint_version': source_blueprint_version,
        'source_stage3_notes_policy_version': source_notes_policy_version,
        'slide_manifest': slide_manifest,
    }
    checkpoint_mgr.save('stage4_pptx', filename, result)
    return result

    candidates = []

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            continue
        if not hasattr(ph, 'text_frame'):
            continue

        left = int(getattr(ph, 'left', 0) or 0)
        width = int(getattr(ph, 'width', 0) or 0)
        height = int(getattr(ph, 'height', 0) or 0)
        area = width * height
        has_text = bool(_clean_text(getattr(ph, 'text', '')))
        ph_type = int(ph.placeholder_format.type)

        candidates.append({
            'idx': idx,
            'left': left,
            'area': area,
            'has_text': has_text,
            'is_body': ph_type == 2,
        })

    if not candidates:
        return None

    min_left = min(c['left'] for c in candidates)
    left_band = [c for c in candidates if c['left'] <= min_left + left_tolerance]
    pool = left_band if left_band else candidates

    pool.sort(
        key=lambda c: (
            1 if c['has_text'] else 0,
            c['area'],
            1 if c['is_body'] else 0,
            -c['left'],
        ),
        reverse=True,
    )
    return pool[0]['idx']


def _remove_slide(prs, slide_index):
    """Remove slide by index using python-pptx internals for dynamic final deck length."""
    sld_id_lst = prs.slides._sldIdLst
    sld_id = sld_id_lst[slide_index]
    rel_id = sld_id.rId
    try:
        prs.part.drop_rel(rel_id)
    except KeyError:
        # Relationship already absent (e.g., slide had embedded media); safe to ignore.
        print(f'    warn: rel {rel_id} not found when removing slide {slide_index}, skipping.')
    del sld_id_lst[slide_index]


def _normalize_blueprint_entry(item, index):
    slide_type = _clean_text(item.get('type', 'content')).lower() or 'content'
    if slide_type not in {'title', 'agenda', 'content', 'image_content', 'diagram', 'conclusion'}:
        slide_type = 'content'

    title = _trim_title(item.get('title') or f'Slide {index}')
    bullets = _trim_bullets(item.get('bullets', []))
    notes = trim_notes(_clean_text(item.get('speaker_notes') or item.get('notes', '')))

    return {
        'type': slide_type,
        'title': title,
        'bullets': bullets,
        'notes': notes,
        'source_slide_num': item.get('source_slide_num'),
        'title_inferred': bool(item.get('title_inferred', False)),
        'embed_source_image': bool(item.get('embed_source_image', False)),
        'image_path': item.get('image_path') or '',
    }


def _ordered_from_typed_blueprint(stage3):
    typed = stage3.get('typed_blueprint')
    if not isinstance(typed, list) or not typed:
        return None

    ordered = []
    for idx, item in enumerate(typed, start=1):
        if not isinstance(item, dict):
            continue
        ordered.append(_normalize_blueprint_entry(item, idx))
    return ordered if ordered else None


def _ordered_from_legacy(parsed, stage3):
    slides_data = parsed['slides']
    insert_order = stage3.get('insert_order', [])
    slide_bullets = stage3.get('slide_bullets', {})
    speaker_notes = stage3.get('speaker_notes', {})
    missing_by_type = {m['content']['slide_type']: m['content'] for m in insert_order if isinstance(m, dict) and m.get('content')}

    ordered = []
    if 'title' in missing_by_type:
        tc = missing_by_type['title']
        ordered.append(_normalize_blueprint_entry({
            'type': 'title',
            'title': tc.get('title'),
            'bullets': tc.get('bullets', []),
            'speaker_notes': tc.get('speaker_notes', ''),
        }, len(ordered) + 1))

    if 'agenda' in missing_by_type:
        ac = missing_by_type['agenda']
        ordered.append(_normalize_blueprint_entry({
            'type': 'agenda',
            'title': ac.get('title', 'Agenda'),
            'bullets': ac.get('bullets', []),
            'speaker_notes': ac.get('speaker_notes', ''),
        }, len(ordered) + 1))

    for s in slides_data:
        slide_num = s['slide_num']
        title_text = s.get('title') or f'Slide {slide_num}'
        raw_text = s.get('raw_text', '')
        notes_raw = speaker_notes.get(str(slide_num), '')

        raw_lines = [ln.strip() for ln in raw_text.split('\n') if ln.strip()]
        if raw_lines and _clean_text(raw_lines[0]).lower() == _clean_text(title_text).lower():
            raw_lines = raw_lines[1:]

        ai_body = slide_bullets.get(str(slide_num), [])
        body = ai_body if isinstance(ai_body, list) and ai_body else clean_body_lines(raw_lines)
        archetype = 'image_content' if s.get('image_path') and int(s.get('word_count', 0) or 0) <= 24 else 'content'

        ordered.append(_normalize_blueprint_entry({
            'type': archetype,
            'source_slide_num': slide_num,
            'title': title_text,
            'bullets': body,
            'speaker_notes': notes_raw,
        }, len(ordered) + 1))

    if 'conclusion' in missing_by_type:
        cc = missing_by_type['conclusion']
        ordered.append(_normalize_blueprint_entry({
            'type': 'conclusion',
            'title': cc.get('title', 'Conclusion'),
            'bullets': cc.get('bullets', []),
            'speaker_notes': cc.get('speaker_notes', ''),
        }, len(ordered) + 1))

    return ordered


def _build_ordered_manifest(parsed, stage3):
    typed = _ordered_from_typed_blueprint(stage3)
    if typed:
        return typed
    return _ordered_from_legacy(parsed, stage3)


# ── Main builder ──────────────────────────────────────────────────────

def build_pptx(filename):
    """Stage 4: build themed PPTX from stage3 manifest using reference.pptx."""

    parsed = checkpoint_mgr.load('stage1_parsed', filename)
    stage3 = checkpoint_mgr.load('stage3_content', filename)

    if not parsed:
        raise Exception('Stage 1 checkpoint missing.')
    if not stage3:
        raise Exception('Stage 3 checkpoint missing.')

    source_blueprint_version = int(stage3.get('typed_blueprint_version', 0) or 0)
    source_notes_policy_version = int(stage3.get('notes_policy_version', 0) or 0)

    if is_cache_reuse_enabled() and checkpoint_mgr.exists('stage4_pptx', filename):
        cached = checkpoint_mgr.load('stage4_pptx', filename)
        if (
            cached
            and 'error' not in cached
            and cached.get('build_policy_version') == BUILD_POLICY_VERSION
            and int(cached.get('source_stage3_blueprint_version', 0) or 0) == source_blueprint_version
            and int(cached.get('source_stage3_notes_policy_version', 0) or 0) == source_notes_policy_version
        ):
            print(f'Valid stage 4 checkpoint for {filename}')
            return cached

    ordered = _build_ordered_manifest(parsed, stage3)

    print(f'Building {len(ordered)} slides...')

    # ── Open reference and modify in-place ───────────────────────────
    if not os.path.exists(REFERENCE_PATH):
        raise FileNotFoundError(
            f'Reference template not found: {REFERENCE_PATH}. '
            'Expected file at ppt-pipeline/theme/reference.pptx'
        )
    prs = Presentation(REFERENCE_PATH)
    ref_slides = list(prs.slides)
    print(f'Reference: {len(ref_slides)} slides')

    # The reference deck is a style scaffold, not a fixed output length.
    # Trim extra scaffold slides so final output follows generated manifest length.
    if len(ref_slides) > len(ordered):
        for idx in range(len(ref_slides) - 1, len(ordered) - 1, -1):
            _remove_slide(prs, idx)
        ref_slides = list(prs.slides)
        print(f'Trimmed reference scaffold to {len(ref_slides)} slides for dynamic output length')

    slide_manifest = []

    for i, item in enumerate(ordered):
        title = item['title']
        bullets = item.get('bullets', [])
        notes = item.get('notes', '')
        slide_type = item.get('type', 'content')

        if i < len(ref_slides):
            slide = ref_slides[i]
        else:
            if slide_type == 'title':
                layout = _find_layout(prs, ['title', 'cover'])
            elif slide_type == 'diagram':
                layout = _find_layout(prs, ['SECTION_TITLE_AND_DESCRIPTION', 'section', 'TITLE_AND_BODY'])
            else:
                layout = _find_layout(prs, ['TITLE AND CONTENT', 'content', 'TITLE_AND_BODY'])
            slide = prs.slides.add_slide(layout)

        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        print(f'  [{i+1}] {slide_type} "{title[:50]}" | phs={ph_indices} | {len(bullets)} bullets')

        # Always strip stale reference template images before writing our own content.
        _remove_content_area_pictures(slide, slide_type)

        # Insert source image for diagram/image_content slides that have a source PNG.
        embed_image = bool(item.get('embed_source_image', False))
        image_path = item.get('image_path') or ''
        if embed_image and image_path:
            _insert_source_image(slide, image_path, slide_type)

        effective_bullets = _trim_bullets(bullets)

        # Apply content to slide
        _set_ph_text(slide, 0, title)

        if slide_type == 'diagram':
            # Diagram slides: the body area is occupied by the infographic image.
            # Writing text bullets here causes overlap. Clear placeholders and
            # leave the body empty — the diagram speaks for itself visually.
            # Speaker notes (narration) still carry the full explanation.
            _set_ph_bullets(slide, 1, [''])
            _set_ph_bullets(slide, 2, [''])
        elif slide_type == 'image_content' and not effective_bullets:
            # image_content with no bullets: showing an empty dashed box looks broken.
            # Downgrade to diagram treatment — clear body, image fills the space.
            print(f'    [{i+1}] image_content has no bullets → treating as diagram (no body text)')
            _set_ph_bullets(slide, 1, [''])
            _set_ph_bullets(slide, 2, [''])
        else:
            # Clear both body placeholders first to avoid reference bleed-through.
            _set_ph_bullets(slide, 1, [''])
            _set_ph_bullets(slide, 2, [''])

            target_body_idx = _resolve_body_placeholder_idx(slide, slide_type)
            if target_body_idx is not None:
                _set_ph_bullets(slide, target_body_idx, effective_bullets if effective_bullets else [''])

        _set_notes(slide, notes)

        slide_manifest.append({
            'slide_num': i + 1,
            'type': slide_type,
            'title': title,
            'bullets': effective_bullets,
            'notes': notes,
            'source_slide_num': item.get('source_slide_num'),
            'title_inferred': bool(item.get('title_inferred', False)),
            'embed_source_image': bool(item.get('embed_source_image', False)),
            'image_path': item.get('image_path') or '',
        })

    # Save the PPTX
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f'{filename}.pptx')
    prs.save(out_path)
    print(f'Saved PPTX to: {out_path}')

    # Calculate totals
    total_words = sum(len(item['notes'].split()) for item in ordered if item.get('notes'))
    est_mins = total_words / 130

    print(f'\nNarrator: ~{total_words} words -> ~{est_mins:.1f} min at 130 wpm')
    final_slide_count = len(prs.slides)
    print(f'Stage 4 complete: {final_slide_count} slides -> {out_path}')

    result = {
        'filename': filename,
        'output_path': out_path,
        'total_slides': final_slide_count,
        'narrator_words': total_words,
        'estimated_minutes': round(est_mins, 2),
        'build_policy_version': BUILD_POLICY_VERSION,
        'source_stage3_blueprint_version': source_blueprint_version,
        'source_stage3_notes_policy_version': source_notes_policy_version,
        'slide_manifest': slide_manifest,
    }
    checkpoint_mgr.save('stage4_pptx', filename, result)
    return result
